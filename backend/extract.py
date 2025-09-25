# backend/extract.py
import io
import pandas as pd
from PyPDF2 import PdfReader
from docx import Document
import docx2txt
from pptx import Presentation

def extract_text(filename: str, file_contents: bytes) -> str:
    extension = filename.split(".")[-1].lower()
    try:
        # ----- PDF -----
        if extension == "pdf":
            try:
                reader = PdfReader(io.BytesIO(file_contents))
                pages = [p.extract_text() for p in reader.pages if p.extract_text()]
                return "\n".join(pages).strip() if pages else "❌ Could not extract text from PDF."
            except Exception:
                return "❌ Could not extract text from PDF."

        # ----- DOCX -----
        if extension == "docx":
            try:
                doc = Document(io.BytesIO(file_contents))
                text = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
                if text.strip():
                    return text.strip()
                # fallback
                return docx2txt.process(io.BytesIO(file_contents))
            except Exception:
                return "❌ Could not extract text from DOCX."

        # ----- TXT -----
        if extension == "txt":
            return file_contents.decode("utf-8", errors="ignore").strip() or "❌ TXT file empty."

        # ----- PPTX -----
        if extension == "pptx":
            try:
                prs = Presentation(io.BytesIO(file_contents))
                slides_text = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slides_text.append(shape.text)
                return "\n".join(slides_text).strip() if slides_text else "❌ Could not extract text from PPTX."
            except Exception:
                return "❌ Could not extract text from PPTX."

        # ----- Excel (XLSX / XLS) -----
        if extension in ["xlsx", "xls"]:
            try:
                excel_data = pd.read_excel(io.BytesIO(file_contents), sheet_name=None)
                text = []
                for sheet_name, df in excel_data.items():
                    text.append(f"Sheet: {sheet_name}")
                    text.append(df.astype(str).apply(lambda x: " ".join(x), axis=1).str.cat(sep="\n"))
                return "\n".join(text).strip()
            except Exception:
                return "❌ Could not extract text from Excel file."

        # ----- CSV -----
        if extension == "csv":
            try:
                df = pd.read_csv(io.BytesIO(file_contents))
                return df.astype(str).apply(lambda x: " ".join(x), axis=1).str.cat(sep="\n") or "❌ CSV file empty."
            except Exception:
                return "❌ Could not extract text from CSV."

        # ----- Unsupported -----
        return f"⚠️ Sorry, {extension.upper()} files are not fully supported yet."

    except Exception as e:
        print(f"Extraction failed for {filename}: {e}")
        return "❌ Could not extract text from this file."
